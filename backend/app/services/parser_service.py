import os
import zipfile
import logging
from PIL import Image
import pypdf
import docx
import pptx
try:
    import pytesseract
    # Try default path for Windows Tesseract if not on PATH
    if os.name == 'nt' and not any(os.path.exists(os.path.join(p, 'tesseract.exe')) for p in os.environ.get('PATH', '').split(os.pathsep)):
        default_win_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(default_win_path):
            pytesseract.pytesseract.tesseract_cmd = default_win_path
except ImportError:
    pytesseract = None

logger = logging.getLogger("memoryverse.parser")

class DocumentParser:
    @staticmethod
    def parse_file(file_path: str) -> str:
        """
        Determines the file type and calls the appropriate parser.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found at {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == ".pdf":
                return DocumentParser.parse_pdf(file_path)
            elif ext in [".docx", ".doc"]:
                return DocumentParser.parse_docx(file_path)
            elif ext in [".pptx", ".ppt"]:
                return DocumentParser.parse_pptx(file_path)
            elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"]:
                return DocumentParser.parse_image(file_path)
            elif ext == ".zip":
                return DocumentParser.parse_zip(file_path)
            elif ext in [".txt", ".md", ".json", ".csv"]:
                return DocumentParser.parse_text_file(file_path)
            else:
                logger.warning(f"Unsupported file extension {ext}. Attempting text reading.")
                return DocumentParser.parse_text_file(file_path)
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {str(e)}")
            # Return basic metadata if extraction fails completely
            return f"Filename: {os.path.basename(file_path)}\nParsing failed with error: {str(e)}"

    @staticmethod
    def parse_pdf(file_path: str) -> str:
        text = ""
        try:
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logger.error(f"PDF text extraction failed: {str(e)}")
        
        # Fallback to OCR if no text was extracted (scanned PDF)
        if not text.strip() and pytesseract:
            logger.info("PDF text extraction empty. Attempting OCR fallback (simulated or real).")
            # In a real environment, we'd convert pages to images and OCR.
            # We'll log it and return basic info.
            return f"Scanned PDF Document: {os.path.basename(file_path)}"
        
        return text

    @staticmethod
    def parse_docx(file_path: str) -> str:
        text = ""
        try:
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    text += " | ".join(cell.text for cell in row.cells) + "\n"
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
        return text

    @staticmethod
    def parse_pptx(file_path: str) -> str:
        text = ""
        try:
            prs = pptx.Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_idx+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
        return text

    @staticmethod
    def parse_image(file_path: str) -> str:
        if pytesseract is None:
            logger.warning("pytesseract is not imported. OCR unavailable. Returning image details.")
            return f"Image File: {os.path.basename(file_path)} (OCR unavailable)"
        
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return text
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {str(e)}")
            return f"Image File: {os.path.basename(file_path)} (OCR failed: {str(e)})"

    @staticmethod
    def parse_zip(file_path: str) -> str:
        extracted_text = []
        temp_dir = os.path.join(os.path.dirname(file_path), "temp_zip_extract")
        os.makedirs(temp_dir, exist_ok=True)
        
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        sub_file_path = os.path.join(root, file)
                        sub_ext = os.path.splitext(file)[1].lower()
                        if sub_ext in [".pdf", ".docx", ".pptx", ".txt", ".png", ".jpg", ".jpeg"]:
                            content = DocumentParser.parse_file(sub_file_path)
                            extracted_text.append(f"=== File: {file} ===\n{content}\n")
                        # Clean up sub file
                        try:
                            os.remove(sub_file_path)
                        except:
                            pass
        except Exception as e:
            logger.error(f"ZIP extraction failed: {str(e)}")
        finally:
            # Clean up temp directory
            try:
                os.rmdir(temp_dir)
            except:
                pass
        
        return "\n".join(extracted_text)

    @staticmethod
    def parse_text_file(file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Text file reading failed: {str(e)}")
            return ""
